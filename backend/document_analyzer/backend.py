#backend.py
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import os
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
import random
import string
import json
from quiz_generator import export_quiz
from flash_card_generator import export_flashcards
from summary_generator import export_summary
from gemini_call import prompt_everyting
import google.generativeai as genai
from speech_to_text import AudioProcessor
import tempfile
import shutil
app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}}, supports_credentials=True)

# Initialize Gemini for chatbot
class Chatbot:
    def __init__(self):
        self.api_key = "AIzaSyDSRv_0xjtnd92cCnvuCAv7QB-1PJcVU1Y"
        genai.configure(api_key=self.api_key)
        self.model = genai.GenerativeModel("gemini-2.0-flash")
        self.chat = None
        self.content = None
        
    def start_chat(self, content: str) -> str:
        """Initialize a new chat session with content"""
        try:
            system_prompt = f"""You are a helpful assistant that answers questions based only on the provided content. 
            If the answer cannot be found in the content, politely say so. Here is the content:

            {content}"""
            
            self.chat = self.model.start_chat(history=[])
            self.chat.send_message(system_prompt)
            self.content = content
            return "Chat started successfully. You can now ask questions about the content."
        except Exception as e:
            return f"Error starting chat: {str(e)}"
    
    def get_response(self, message: str) -> str:
        """Get response for user message"""
        try:
            if not self.chat:
                return "Please send content first to start the chat."
            
            response = self.chat.send_message(message)
            return response.text
        except Exception as e:
            return f"Error generating response: {str(e)}"
    
    def clear_chat(self) -> str:
        """Clear current chat session"""
        self.chat = None
        self.content = None
        return "Chat cleared successfully."

# Create global chatbot instance
chatbot = Chatbot()

def generate_quiz(flashcards):
    quiz = []
    for card in flashcards:
        if random.choice([True, False]):
            question = {"question": card[1], "possible_answers": [], "index": -1}
            incorrect_answers = []
            for other_card in flashcards:
                if other_card != card:
                    incorrect_answers.append(other_card[0])
            question["possible_answers"] = random.sample(incorrect_answers, 3)
            question["index"] = random.randint(0, 3)
            question["possible_answers"].insert(question["index"], card[0])
        else:
            question = {"question": card[0], "possible_answers": [], "index": -1}
            incorrect_answers = []
            for other_card in flashcards:
                if other_card != card:
                    incorrect_answers.append(other_card[1])
            question["possible_answers"] = random.sample(incorrect_answers, 3)
            question["index"] = random.randint(0, 3)
            question["possible_answers"].insert(question["index"], card[1])
        quiz.append(question)
    return quiz

def generate_id():
    ids = set()
    for file in os.listdir(os.getcwd() + "/database"):
        if file.endswith(".json"):
            ids.add(file)

    letters = string.ascii_letters + string.digits
    id = "".join(random.choice(letters) for _ in range(5))
    while id in ids:
        id = "".join(random.choice(letters) for _ in range(5))
    return id

def handle_pdf(pdf):
    reader = PdfReader(os.getcwd() + "/file.pdf")
    number_of_pages = len(reader.pages)
    s = ""
    for i in range(number_of_pages):
        page = reader.pages[i]
        text = page.extract_text()
        s += text
    return s

def handle_txt(txt):
    s = open("file.txt").read()
    return s

def handle_docx(docx):
    d = Document(os.getcwd() + "/file.docx")
    s = ""
    for paragraph in d.paragraphs:
        s += paragraph.text + "\n"
    return s

def handle_pptx(pptx):
    p = Presentation("file.pptx")
    s = ""
    for slide in p.slides:
        try:
            for shape in slide.shapes:
                try:
                    s += shape.text
                except Exception:
                    pass
            s += "\n"
        except Exception:
            pass
    return s


#-----------------------------------------

# Add these imports at the top of backend.py
from goal_todo_features import GoalTodoManager
from video_search import VideoSearchManager

# Add these class initializations after the chatbot initialization
goal_todo_manager = GoalTodoManager(api_key="AIzaSyDSRv_0xjtnd92cCnvuCAv7QB-1PJcVU1Y")  # Use the same API key as chatbot
video_search_manager = VideoSearchManager()

# Add these new routes to backend.py
@app.route("/goaladvise", methods=["POST"])
def goal_advise():
    try:
        data = request.json
        if not data or 'goal' not in data:
            return jsonify({"error": "No goal statement provided"}), 400
        
        result = goal_todo_manager.get_goal_advice(data['goal'])
        return jsonify(result), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/generatetodotask", methods=["POST"])
def generate_todo_auto():
    try:
        data = request.json
        if not data or 'goal' not in data:
            return jsonify({"error": "No goal provided"}), 400
        
        result = goal_todo_manager.generate_todo_tasks(data['goal'])
        return jsonify(result), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/search_videos", methods=["POST"])
def search_videos():
    try:
        data = request.json
        if not data or 'query' not in data:
            return jsonify({"error": "No search query provided"}), 400
        
        result = video_search_manager.format_video_response(data['query'])
        return jsonify(result), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# Routes
@app.route("/")
def hello_world():
    return "Hello, World!"


# @app.route("/upload", methods=["POST"])
# def upload():
#     try:
#         file = request.files["file"]
#         if not file:
#             return jsonify({"error": "No file provided"}), 400
            
#         name = file.filename
#         extension = name.split(".")[-1].lower()
        
#         if extension not in ["pdf", "txt", "docx", "pptx"]:
#             return jsonify({"error": "Unsupported file format"}), 400
            
#         file.save(os.path.join(os.getcwd(), f"file.{extension}"))
#         file.close()

#         try:
#             if extension == "pdf":
#                 s = handle_pdf(file)
#             elif extension == "txt":
#                 s = handle_txt(file)
#             elif extension == "docx":
#                 s = handle_docx(file)
#             elif extension == "pptx":
#                 s = handle_pptx(file)
#             else:
#                 return jsonify({"error": "Unsupported file format"}), 400

#             if not s:
#                 return jsonify({"error": "Failed to extract text from file"}), 400

#             response = prompt_everyting(s)
#             if not response:
#                 return jsonify({"error": "Failed to process content"}), 500

#             quiz2 = generate_quiz(response["flash_cards"])
#             response["quiz"] += quiz2
#             random.shuffle(response["quiz"])

#             id = generate_id()
#             response["id"] = id
            
#             # Initialize chat with the document content
#             chatbot.start_chat(s)
            
#             # Save to database
#             save_path = os.path.join(os.getcwd(), "database", f"{id}.json")
#             os.makedirs(os.path.dirname(save_path), exist_ok=True)
            
#             with open(save_path, "w", encoding='utf-8') as f:
#                 json.dump(response, f, ensure_ascii=False, indent=2)

#             return jsonify(response), 200

#         except Exception as e:
#             print(f"Processing error: {str(e)}")
#             return jsonify({"error": f"Failed to process file: {str(e)}"}), 500
            
#     except Exception as e:
#         print(f"Upload error: {str(e)}")
#         return jsonify({"error": f"Upload failed: {str(e)}"}), 500

def handle_pdf(file_path):
    reader = PdfReader(file_path)
    number_of_pages = len(reader.pages)
    s = ""
    for i in range(number_of_pages):
        page = reader.pages[i]
        text = page.extract_text()
        s += text
    return s

def handle_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def handle_docx(file_path):
    d = Document(file_path)
    s = ""
    for paragraph in d.paragraphs:
        s += paragraph.text + "\n"
    return s

def handle_pptx(file_path):
    p = Presentation(file_path)
    s = ""
    for slide in p.slides:
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    s += shape.text + "\n"
        except Exception:
            continue
    return s

@app.route("/upload", methods=["POST"])
def upload():
    temp_dir = None
    try:
        file = request.files["file"]
        if not file:
            return jsonify({"error": "No file provided"}), 400
            
        name = file.filename
        extension = name.split(".")[-1].lower()
        
        # Create a temporary directory for file processing
        temp_dir = tempfile.mkdtemp()
        file_path = os.path.join(temp_dir, f"temp_file.{extension}")
        
        # Save the uploaded file
        file.save(file_path)
        
        try:
            # Initialize text content
            text_content = ""
            
            # Handle different file types
            if extension == "pdf":
                text_content = handle_pdf(file_path)
            elif extension == "txt":
                text_content = handle_txt(file_path)
            elif extension == "docx":
                text_content = handle_docx(file_path)
            elif extension == "pptx":
                text_content = handle_pptx(file_path)
            elif extension in ['mp3', 'wav', 'm4a', 'ogg', 'mp4', 'avi', 'mov', 'mkv']:
                processor = AudioProcessor()
                text_content = processor.process_file(file_path)
            else:
                return jsonify({"error": "Unsupported file format"}), 400

            if not text_content:
                return jsonify({"error": "Failed to extract text from file"}), 400

            # Process the text content
            response = prompt_everyting(text_content)
            if not response:
                return jsonify({"error": "Failed to process content"}), 500

            # Generate additional quiz questions and shuffle
            quiz2 = generate_quiz(response["flash_cards"])
            response["quiz"] += quiz2
            random.shuffle(response["quiz"])

            # Generate unique ID and save result
            id = generate_id()
            response["id"] = id
            
            # Initialize chat with the document content
            chatbot.start_chat(text_content)
            
            # Save to database
            save_path = os.path.join(os.getcwd(), "database", f"{id}.json")
            os.makedirs(os.path.dirname(save_path), exist_ok=True)
            
            with open(save_path, "w", encoding='utf-8') as f:
                json.dump(response, f, ensure_ascii=False, indent=2)

            return jsonify(response), 200

        except Exception as e:
            print(f"Processing error: {str(e)}")
            return jsonify({"error": f"Failed to process file: {str(e)}"}), 500
            
    except Exception as e:
        print(f"Upload error: {str(e)}")
        return jsonify({"error": f"Upload failed: {str(e)}"}), 500
    
    finally:
        # Clean up temporary directory and files
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except Exception as e:
                print(f"Warning: Failed to clean up temporary files: {str(e)}")

@app.route("/fetch_id", methods=["POST"])
def fetch_id():
    id = request.json.get("id")
    try:
        with open(os.getcwd() + "/database/" + id + ".json", "r") as f:
            data = json.load(f)
        return data
    except:
        return {"summary": 404}

# @app.route("/recent", methods=["GET"])
# def recent():
#     output = {"recent": []}
#     i = 0
#     for file in os.listdir(os.getcwd() + "/database"):
#         if i == 10:
#             break
#         if file.endswith(".json"):
#             with open(os.getcwd() + "/database/" + file, "r") as f:
#                 data = json.load(f)
#                 id = file.split(".json")[0]
#                 title = data["title"]
#                 output["recent"].append({"id": id, "title": title})
#         i += 1
#     return output

@app.route("/recent", methods=["GET"])
def recent():
    output = {"recent": []}
    i = 0
    for file in os.listdir(os.getcwd() + "/database"):
        if i == 10:
            break
        if file.endswith(".json"):
            try:
                # Fix: Add UTF-8 encoding when opening the file
                with open(os.getcwd() + "/database/" + file, "r", encoding='utf-8') as f:
                    data = json.load(f)
                    id = file.split(".json")[0]
                    title = data["title"]
                    output["recent"].append({"id": id, "title": title})
                i += 1
            except Exception as e:
                print(f"Error reading file {file}: {str(e)}")
                continue
    return output

@app.route("/export", methods=["POST"])
def export():
    selected = request.json.get("selected")
    data = request.json.get("data")
    if selected == 0:
        data = data["summary"]
        export_summary(data, "Summary.docx")
        return send_file("Summary.docx", as_attachment=True)
    elif selected == 1:
        data = data["flash_cards"]
        export_flashcards(data, "Flashcards.docx")
        return send_file("Flashcards.docx", as_attachment=True)
    else:
        data = data["quiz"]
        export_quiz(data, "Quiz.docx")
        return send_file("Quiz.docx", as_attachment=True)

@app.route("/chat/start", methods=["POST"])
def start_chat():
    try:
        data = request.json
        content = data.get("content")
        
        if not content:
            return jsonify({"error": "No content provided"}), 400
            
        result = chatbot.start_chat(content)
        return jsonify({"message": result})
    except Exception as e:
        return jsonify({"error": str(e)}), 400

@app.route("/chat/message", methods=["POST"])
def chat_message():
    try:
        data = request.json
        message = data.get("message")
        
        if not message:
            return jsonify({"error": "No message provided"}), 400
        
        response = chatbot.get_response(message)
        return jsonify({"response": response})
    except Exception as e:
        return jsonify({"error": str(e)}), 400

@app.route("/chat/clear", methods=["POST"])
def clear_chat():
    try:
        result = chatbot.clear_chat()
        return jsonify({"message": result})
    except Exception as e:
        return jsonify({"error": str(e)}), 400

if __name__ == "__main__":
    app.run(debug=True, port=5000)
    
